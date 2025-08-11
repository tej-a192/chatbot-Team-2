# server/rag_service/document_generator.py
import re
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DocxInches
import logging
from prompts import (
    PPTX_GENERATION_FROM_TOPIC_PROMPT_TEMPLATE,
    DOCX_GENERATION_FROM_TOPIC_PROMPT_TEMPLATE
)

logger = logging.getLogger(__name__)

# --- PROMPT FOR INTELLIGENT PPTX GENERATION (JSON-based) ---
PPTX_EXPANSION_PROMPT_TEMPLATE = """
You are a professional presentation designer and subject matter expert.
Your task is to expand a given OUTLINE (which could be a list of key topics or FAQs) into a full, detailed, 6-8 slide presentation.
You must use the provided SOURCE DOCUMENT TEXT as your only source of truth. Do not use outside knowledge.
Your output MUST be a single, valid JSON array, where each object represents a slide.

**JSON Object Schema for each slide:**
{{
  "slide_title": "A concise and engaging title for the slide.",
  "slide_content": "Detailed, professional paragraph(s) and/or bullet points elaborating on the outline point. This text will be displayed on the slide. Use Markdown for formatting (e.g., **bold**, *italics*, - bullet points).",
  "image_prompt": "A highly descriptive, creative prompt for an AI text-to-image model (like DALL-E or Midjourney) to generate a relevant and visually appealing image for this specific slide. Describe the style, subject, and composition. Example: 'A photorealistic image of a futuristic server room with glowing blue data streams flowing between racks, symbolizing data processing. Cinematic lighting.'"
}}

**INSTRUCTIONS:**
1.  **Analyze Outline & Source:** For each point in the OUTLINE, create at least one slide object in the JSON array.
2.  **Expand Content:** Elaborate on each outline point using only information from the SOURCE DOCUMENT TEXT.
3.  **Create Image Prompts:** For each slide, generate a unique and descriptive `image_prompt` that visually represents the slide's content.
4.  **JSON Format:** Ensure the final output is a single, clean JSON array with no other text before or after it.

---
**SOURCE DOCUMENT TEXT (Your knowledge base):**
{source_document_text}
---
**OUTLINE (Topics/FAQs to expand into a presentation):**
{outline_content}
---

**FINAL PRESENTATION JSON ARRAY:**
"""

# --- PROMPT FOR INTELLIGENT DOCX GENERATION (Markdown-based) ---
DOCX_EXPANSION_PROMPT_TEMPLATE = """
You are a professional content creator and subject matter expert.
Your task is to expand a given OUTLINE (which could be a list of key topics or FAQs) into a full, detailed, multi-page document in Markdown format.
You must use the provided SOURCE DOCUMENT TEXT as your only source of truth. Do not use outside knowledge.
The final output must be a single block of well-structured Markdown text.

**INSTRUCTIONS:**
1.  **Main Title:** Start the document with a main title using H1 syntax (e.g., `# Expanded Report on Key Topics`).
2.  **Section per Outline Point:** For each point in the OUTLINE, create a detailed section with a clear H2 or H3 heading (e.g., `## Topic Name`).
3.  **Content Expansion:** For each section, write detailed, professional paragraphs that elaborate on the outline point. Extract relevant facts, figures, and explanations from the SOURCE DOCUMENT TEXT.
4.  **Markdown Usage:** Use bullet points, bold text, and clear paragraphs to structure the content effectively.

---
**SOURCE DOCUMENT TEXT (Your knowledge base):**
{source_document_text}
---
**OUTLINE (Topics/FAQs to expand into a document):**
{outline_content}
---

**FINAL DOCUMENT MARKDOWN:**
"""

def expand_content_with_llm(outline_content, source_document_text, doc_type, llm_function):
    """Uses an LLM to expand an outline into full content for the specified doc type."""
    logger.info(f"Expanding outline for '{doc_type}' using LLM...")
    
    if doc_type == 'pptx':
        prompt = PPTX_EXPANSION_PROMPT_TEMPLATE.format(
            source_document_text=source_document_text,
            outline_content=outline_content
        )
    else: # for 'docx'
        prompt = DOCX_EXPANSION_PROMPT_TEMPLATE.format(
            source_document_text=source_document_text,
            outline_content=outline_content
        )

    expanded_content = llm_function(prompt)
    
    if not expanded_content or not expanded_content.strip():
        raise ValueError("LLM failed to generate expanded content.")
    
    logger.info(f"LLM generated expanded content for {doc_type}. Length: {len(expanded_content)}")
    return expanded_content

def parse_pptx_json(json_string: str) -> list:
    """
    Parses the LLM's JSON output for PPTX generation with enhanced error handling.
    """
    try:
        # First, find the JSON block. This is more robust against preamble/apology text from the LLM.
        json_match = re.search(r'\[\s*\{[\s\S]*?\}\s*\]', json_string, re.DOTALL)
        if not json_match:
            raise ValueError("No valid JSON array of slides was found in the AI's response.")
        
        cleaned_str = json_match.group(0)
        slides_data = json.loads(cleaned_str)
        
        # Now, validate the structure. It must be a list, and not an empty one.
        if not isinstance(slides_data, list) or not slides_data:
            raise ValueError("Parsed JSON is not a non-empty list of slides.")
            
        # Optional: Deeper validation of individual slide objects
        for i, slide in enumerate(slides_data):
            if not isinstance(slide, dict) or "slide_title" not in slide or "slide_content" not in slide:
                raise ValueError(f"Slide object at index {i} is missing required 'slide_title' or 'slide_content' keys.")

        return slides_data
    except (json.JSONDecodeError, ValueError) as e:
        logger.error(f"Failed to parse or validate JSON from LLM response: {e}\nRaw Response Preview: {json_string[:500]}")
        # Return an empty list to signal failure to the calling function in app.py
        return []

def refined_parse_docx_markdown(markdown_content: str) -> list:
    """
    Parses the expanded markdown for DOCX generation. Returns an empty list on failure.
    """
    if not markdown_content or not markdown_content.strip():
        return [] # Return empty list if there's no content
    
    title_match = re.search(r"^\s*#\s+(.*)", markdown_content, re.MULTILINE)
    if title_match:
        title = title_match.group(1).strip()
        content = markdown_content[title_match.end():].strip()
    else:
        # If there's no H1, it might just be paragraphs. Use a default title.
        title = "Generated Document"
        content = markdown_content

    # If after parsing, there's no content left, it's a failure.
    if not content.strip():
        return []

    return [{"title": title, "text_content": content}]


def add_text_to_shape_with_markdown(text_frame, markdown_text, is_title=False, is_notes=False):
    text_frame.clear()
    text_frame.word_wrap = True
    title_font_size = Pt(36)
    content_font_size = Pt(16)
    notes_font_size = Pt(11)

    for line in markdown_text.split('\n'):
        p = text_frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        bullet_match = re.match(r'^(\s*)[\*\-]\s*(.*)', line)
        
        if bullet_match and not is_title:
            leading_spaces, content_line = bullet_match.groups()
            p.level = min(len(leading_spaces) // 2, 5)
        else:
            content_line = line.lstrip()

        segments = re.split(r'(\*\*.*?\*\*|__.*?__)', content_line)
        for segment in segments:
            if not segment: continue
            run = p.add_run()
            if (segment.startswith("**") and segment.endswith("**")) or (segment.startswith("__") and segment.endswith("__")):
                run.text = segment[2:-2]
                run.font.bold = True
            else:
                run.text = segment
            
            if is_title:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.size = title_font_size
            elif is_notes:
                run.font.color.rgb = RGBColor(210, 210, 230)
                run.font.size = notes_font_size
                run.font.italic = True
            else:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.size = content_font_size

def create_ppt(slides_data, output_path):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)

        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1.0), Inches(1.0))
        add_text_to_shape_with_markdown(title_shape.text_frame, slide_data.get("slide_title", "Untitled Slide"), is_title=True)

        content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(8.5), Inches(7.0))
        add_text_to_shape_with_markdown(content_shape.text_frame, slide_data.get("slide_content", "[No content provided]"))

        notes_shape = slide.shapes.add_textbox(Inches(9.5), Inches(1.3), Inches(6.0), Inches(7.0))
        image_prompt_header = "🎨 Image Generation Prompt:"
        image_prompt_body = slide_data.get("image_prompt", "N/A")
        add_text_to_shape_with_markdown(notes_shape.text_frame, f"**{image_prompt_header}**\n{image_prompt_body}", is_notes=True)

    prs.save(output_path)
    return True

def add_markdown_line_to_docx(doc, markdown_line):
    heading_match = re.match(r'^(#+)\s+(.*)', markdown_line)
    if heading_match:
        level = len(heading_match.group(1))
        doc.add_heading(heading_match.group(2).strip(), level=min(level, 4))
        return

    bullet_match = re.match(r'^(\s*)[\*\-]\s+(.*)', markdown_line)
    if bullet_match:
        leading_spaces, content_line = bullet_match.groups()
        p = doc.add_paragraph(style='List Bullet')
        p.paragraph_format.left_indent = DocxInches(0.25 * (len(leading_spaces) // 2))
    else:
        content_line = markdown_line
        p = doc.add_paragraph()
    
    segments = re.split(r'(\*\*.*?\*\*|__.*?__)', content_line)
    for segment in segments:
        if not segment: continue
        run = p.add_run()
        if (segment.startswith("**") and segment.endswith("**")) or (segment.startswith("__") and segment.endswith("__")):
            run.text = segment[2:-2]
            run.font.bold = True
        else:
            run.text = segment

def create_doc(slides_data, output_path, content_key="text_content"):
    doc = Document()
    if slides_data:
        doc_title = slides_data[0].get("title", "Generated Document")
        doc.add_heading(doc_title, level=0)
        
        content_to_add = slides_data[0].get(content_key, "")
        if content_to_add.strip():
            for line in content_to_add.split('\n'):
                add_markdown_line_to_docx(doc, line)
    else:
        doc.add_paragraph("[No content to generate]")
    doc.save(output_path)
    return True


def generate_content_from_topic(topic, doc_type, llm_function):
    """Uses an LLM to generate document content from scratch based on a topic."""
    logger.info(f"Generating content for a new '{doc_type}' on topic: '{topic}'")

    if doc_type == 'pptx':
        prompt = PPTX_GENERATION_FROM_TOPIC_PROMPT_TEMPLATE.format(topic=topic)
    else: # for 'docx'
        prompt = DOCX_GENERATION_FROM_TOPIC_PROMPT_TEMPLATE.format(topic=topic)

    generated_content = llm_function(prompt)
    
    if not generated_content or not generated_content.strip():
        raise ValueError("LLM failed to generate content from the topic.")
    
    logger.info(f"LLM generated content from topic. Length: {len(generated_content)}")
    return generated_content