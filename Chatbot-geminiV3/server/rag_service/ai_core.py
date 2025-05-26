# ./ai_core.py

# Standard Library Imports
import os
import io
import re
import logging
import copy # Used in chunk_document_into_segments
import uuid # For Qdrant Point IDs

# Typing Imports
from typing import Optional, List, Dict, Any, Callable

# --- Configuration Import ---
# Assumes 'server/config.py' is the actual config file.
# `import config` will work if 'server/' directory is in sys.path.
try:
    import config # This should import server/config.py
except ImportError as e:
    logging.basicConfig(level=logging.CRITICAL) 
    logging.getLogger(__name__).critical(
        f"CRITICAL: Failed to import 'config' (expected server/config.py): {e}. "
        "Application will use dummy defaults and likely not function as intended."
    )
    class DummyConfig: # Fallback with CORRECTED attribute names
        DOCUMENT_EMBEDDING_MODEL_NAME = 'sentence-transformers/all-MiniLM-L6-v2' # Default model
        AI_CORE_CHUNK_SIZE = 512           # Corrected
        AI_CORE_CHUNK_OVERLAP = 100        # Corrected
        SPACY_MODEL_NAME = 'en_core_web_sm' # Default SpaCy
        MAX_TEXT_LENGTH_FOR_NER = 500000
    config = DummyConfig()

# --- Global Initializations ---

# 1. Logger Setup
logger = logging.getLogger(__name__)
if not logger.hasHandlers():
    _handler = logging.StreamHandler()
    _formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
    _handler.setFormatter(_formatter)
    logger.addHandler(_handler)
    logger.setLevel(logging.INFO)

# 2. Optional Library Imports and Availability Flags / Placeholders
_AI_CORE_PYPDF_PDFREADERROR = None
try:
    import pypdf
    _AI_CORE_PYPDF_PDFREADERROR = pypdf.errors.PdfReadError
    PYPDF_AVAILABLE = True
except ImportError:
    pypdf = None
    _AI_CORE_PYPDF_PDFREADERROR = type('PdfReadErrorDummy', (Exception,), {})
    PYPDF_AVAILABLE = False
    logger.info("pypdf library not found. PDF parsing with pypdf unavailable.")

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DocxDocument = None
    DOCX_AVAILABLE = False
    logger.info("python-docx library not found. DOCX parsing unavailable.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    PPTX_AVAILABLE = False
    logger.info("python-pptx library not found. PPTX parsing unavailable.")

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    pdfplumber = None
    PDFPLUMBER_AVAILABLE = False
    logger.info("pdfplumber library not found. Rich PDF extraction unavailable.")

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    pd = None
    PANDAS_AVAILABLE = False
    logger.warning("pandas library not found. Table processing might be impacted.")

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    Image = None
    PIL_AVAILABLE = False
    logger.info("Pillow (PIL) library not found. Image processing unavailable.")

try:
    import fitz  # PyMuPDF
    FITZ_AVAILABLE = True
except ImportError:
    fitz = None
    FITZ_AVAILABLE = False
    logger.info("PyMuPDF (fitz) library not found. PDF image extraction via fitz unavailable.")

_AI_CORE_TESSERACT_NOT_FOUND_ERROR = None 
try:
    import pytesseract
    PYTESSERACT_AVAILABLE = True
    _AI_CORE_TESSERACT_NOT_FOUND_ERROR = pytesseract.TesseractNotFoundError
except ImportError:
    pytesseract = None
    PYTESSERACT_AVAILABLE = False
    logger.info("pytesseract library not found. OCR functionality unavailable.")

try:
    import PyPDF2 
    PYPDF2_AVAILABLE = True
except ImportError:
    PyPDF2 = None
    PYPDF2_AVAILABLE = False
    logger.info("PyPDF2 library not found (used for some PDF metadata).")

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    LANGCHAIN_SPLITTER_AVAILABLE = True
except ImportError:
    try:
        from langchain.text_splitter import RecursiveCharacterTextSplitter # Fallback
        LANGCHAIN_SPLITTER_AVAILABLE = True
        logger.debug("Using RecursiveCharacterTextSplitter from langchain.text_splitter.")
    except ImportError:
        RecursiveCharacterTextSplitter = None
        LANGCHAIN_SPLITTER_AVAILABLE = False
        logger.critical("Langchain text splitter not found. Text chunking will fail.")

# 3. NLP Model Initializations
nlp_spacy_core = None
SPACY_MODEL_LOADED = False
try:
    import spacy
    SPACY_LIB_AVAILABLE = True
except ImportError:
    SPACY_LIB_AVAILABLE = False
    logger.warning("spacy library not found. SpaCy model cannot be loaded.")

if SPACY_LIB_AVAILABLE:
    spacy_model_name_from_config = getattr(config, 'SPACY_MODEL_NAME', None)
    if spacy_model_name_from_config:
        try:
            nlp_spacy_core = spacy.load(spacy_model_name_from_config)
            SPACY_MODEL_LOADED = True
            logger.info(f"SpaCy model '{spacy_model_name_from_config}' loaded successfully.")
        except OSError:
            logger.error(
                f"SpaCy model '{spacy_model_name_from_config}' not found. "
                f"Download with: 'python -m spacy download {spacy_model_name_from_config}'. "
                "NER/lemmatization impacted."
            )
        except Exception as e:
            logger.error(f"Unexpected error loading SpaCy model '{spacy_model_name_from_config}': {e}")
    else:
        logger.warning("config.SPACY_MODEL_NAME not defined. SpaCy model not loaded.")

document_embedding_model = None
EMBEDDING_MODEL_LOADED = False
try:
    from sentence_transformers import SentenceTransformer
    SENTENCE_TRANSFORMERS_LIB_AVAILABLE = True
except ImportError:
    SENTENCE_TRANSFORMERS_LIB_AVAILABLE = False
    logger.critical("sentence-transformers library not found. Embedding generation will fail.")

if SENTENCE_TRANSFORMERS_LIB_AVAILABLE:
    # CORRECTED: Use DOCUMENT_EMBEDDING_MODEL_NAME from server/config.py
    embedding_model_name_from_config = getattr(config, 'DOCUMENT_EMBEDDING_MODEL_NAME', None) 
    if embedding_model_name_from_config:
        try:
            document_embedding_model = SentenceTransformer(embedding_model_name_from_config)
            EMBEDDING_MODEL_LOADED = True
            logger.info(f"Sentence Transformer embedding model '{embedding_model_name_from_config}' loaded successfully.")
        except Exception as e:
            logger.error(f"Error loading Sentence Transformer model '{embedding_model_name_from_config}': {e}. Embedding generation will fail.")
    else:
        # CORRECTED: Log message to reflect the correct config variable name
        logger.critical(
            "config.DOCUMENT_EMBEDDING_MODEL_NAME not defined or empty. " 
            "Sentence Transformer model not loaded. Embedding generation will fail."
        )

# --- Stage 1: File Parsing and Raw Content Extraction --- (Functions as previously corrected)
def _parse_pdf_content(file_path: str) -> Optional[str]:
    if not PYPDF_AVAILABLE or not pypdf:
        logger.error("pypdf library not available. PDF parsing with pypdf will fail.")
        return None
    text_content = ""
    try:
        reader = pypdf.PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text()
                if page_text: text_content += page_text + "\n"
            except Exception as page_err:
                logger.warning(f"pypdf: Error extracting text from page {i+1} of {os.path.basename(file_path)}: {page_err}")
        return text_content.strip() or None
    except FileNotFoundError:
        logger.error(f"pypdf: File not found: {file_path}"); return None
    except _AI_CORE_PYPDF_PDFREADERROR as pdf_err: 
        logger.error(f"pypdf: Error reading PDF {os.path.basename(file_path)}: {pdf_err}"); return None
    except Exception as e:
        logger.error(f"pypdf: Unexpected error parsing PDF {os.path.basename(file_path)}: {e}", exc_info=True); return None

def _parse_docx_content(file_path: str) -> Optional[str]:
    if not DOCX_AVAILABLE or not DocxDocument:
        logger.error("python-docx library not available. DOCX parsing will fail.")
        return None
    try:
        doc = DocxDocument(file_path)
        text_content = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        return text_content.strip() or None
    except FileNotFoundError:
        logger.error(f"docx: File not found: {file_path}"); return None
    except Exception as e:
        logger.error(f"docx: Error parsing DOCX {os.path.basename(file_path)}: {e}", exc_info=True); return None

def _parse_txt_content(file_path: str) -> Optional[str]:
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f: text_content = f.read()
        return text_content.strip() or None
    except FileNotFoundError:
        logger.error(f"txt: File not found: {file_path}"); return None
    except Exception as e:
        logger.error(f"txt: Error parsing TXT {os.path.basename(file_path)}: {e}", exc_info=True); return None

def _parse_pptx_content(file_path: str) -> Optional[str]:
    if not PPTX_AVAILABLE or not Presentation:
        logger.warning(f"python-pptx not available. PPTX parsing for {os.path.basename(file_path)} skipped.")
        return None
    text_content = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip(): text_content += shape.text.strip() + "\n\n"
        return text_content.strip() or None
    except FileNotFoundError:
        logger.error(f"pptx: File not found: {file_path}"); return None
    except Exception as e:
        logger.error(f"pptx: Error parsing PPTX {os.path.basename(file_path)}: {e}", exc_info=True); return None

def _get_parser_for_file(file_path: str) -> Optional[Callable]:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.pdf': return _parse_pdf_content
    if ext == '.docx': return _parse_docx_content
    if ext == '.pptx': return _parse_pptx_content
    if ext in ['.txt', '.py', '.js', '.md', '.log', '.csv', '.html', '.xml', '.json']: return _parse_txt_content
    logger.warning(f"Unsupported file extension for basic parsing: {ext} ({os.path.basename(file_path)})")
    return None

def extract_raw_content_from_file(file_path: str) -> Dict[str, Any]:
    file_base_name = os.path.basename(file_path)
    logger.info(f"Starting raw content extraction for: {file_base_name}")
    text_content, tables, images, is_scanned = "", [], [], False
    file_extension = os.path.splitext(file_path)[1].lower()

    parser_func = _get_parser_for_file(file_path)
    if parser_func:
        initial_text = parser_func(file_path)
        if initial_text: text_content = initial_text

    if file_extension == '.pdf':
        if PDFPLUMBER_AVAILABLE and pdfplumber:
            try:
                with pdfplumber.open(file_path) as pdf:
                    pdfplumber_text_parts = [p.extract_text(x_tolerance=1, y_tolerance=1) or "" for p in pdf.pages]
                    pdfplumber_text = "\n".join(pdfplumber_text_parts)
                    clean_pdfplumber_text_len = len(pdfplumber_text.replace("\n", ""))

                    if len(pdf.pages) > 0 and (clean_pdfplumber_text_len < 50 * len(pdf.pages) and clean_pdfplumber_text_len < 200):
                        is_scanned = True; logger.info(f"PDF {file_base_name} potentially scanned (low text from pdfplumber).")
                    
                    if len(pdfplumber_text.strip()) > len(text_content.strip()): text_content = pdfplumber_text.strip()
                    elif not text_content.strip() and pdfplumber_text.strip(): text_content = pdfplumber_text.strip()
                    
                    for page_num, page in enumerate(pdf.pages): 
                        page_tables_data = page.extract_tables()
                        if page_tables_data:
                            for table_data_list in page_tables_data:
                                if table_data_list and PANDAS_AVAILABLE and pd:
                                    try:
                                        if len(table_data_list) > 1 and all(isinstance(c, str) or c is None for c in table_data_list[0]):
                                            tables.append(pd.DataFrame(table_data_list[1:], columns=table_data_list[0]))
                                        else: tables.append(table_data_list)
                                    except Exception as df_err: logger.warning(f"pdfplumber: DataFrame conversion error for table on page {page_num} of {file_base_name}: {df_err}"); tables.append(table_data_list)
                                elif table_data_list: 
                                    tables.append(table_data_list)
                    if tables: logger.info(f"pdfplumber: Extracted {len(tables)} tables from {file_base_name}.")
            except Exception as e: logger.warning(f"pdfplumber: Error during rich PDF processing for {file_base_name}: {e}", exc_info=True)
        
        elif not text_content.strip() and PYPDF_AVAILABLE and pypdf: 
            try:
                if len(pypdf.PdfReader(file_path).pages) > 0: is_scanned = True; logger.info(f"PDF {file_base_name} potentially scanned (no text from pypdf and pdfplumber not used/failed).")
            except: pass

        if FITZ_AVAILABLE and fitz and PIL_AVAILABLE and Image: 
            try:
                doc = fitz.open(file_path)
                for page_idx in range(len(doc)):
                    for img_info in doc.get_page_images(page_idx):
                        try: images.append(Image.open(io.BytesIO(doc.extract_image(img_info[0])["image"])))
                        except Exception as img_err: logger.warning(f"fitz: Could not open image xref {img_info[0]} from page {page_idx} of {file_base_name}: {img_err}")
                if images: logger.info(f"fitz: Extracted {len(images)} images from {file_base_name}.")
                doc.close()
            except Exception as e: logger.warning(f"fitz: Error extracting images from PDF {file_base_name}: {e}", exc_info=True)

    if not text_content.strip() and file_extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif']:
        is_scanned = True
        if PIL_AVAILABLE and Image:
            try: images.append(Image.open(file_path))
            except Exception as e: logger.error(f"Could not open image file {file_base_name}: {e}", exc_info=True)

    logger.info(f"Raw content extraction complete for {file_base_name}. Text length: {len(text_content)}, Tables: {len(tables)}, Images: {len(images)}, Scanned: {is_scanned}")
    return {'text_content': text_content.strip(), 'tables': tables, 'images': images, 'is_scanned': is_scanned, 'file_type': file_extension}

# --- Stage 2: OCR --- (Function as previously corrected)
def perform_ocr_on_images(image_objects: List[Any]) -> str:
    if not image_objects: return ""
    if not PYTESSERACT_AVAILABLE or not pytesseract:
        logger.error("Pytesseract is not available. OCR cannot be performed.")
        return ""

    logger.info(f"Performing OCR on {len(image_objects)} image(s).")
    ocr_text_parts = []
    images_ocrd = 0
    for i, img in enumerate(image_objects):
        try:
            if not (PIL_AVAILABLE and Image and isinstance(img, Image.Image)):
                logger.warning(f"Skipping non-PIL Image object at index {i} for OCR.")
                continue
            text = pytesseract.image_to_string(img.convert('L'))
            if text and text.strip(): 
                ocr_text_parts.append(text.strip())
                images_ocrd += 1
        except Exception as e:
            if _AI_CORE_TESSERACT_NOT_FOUND_ERROR and isinstance(e, _AI_CORE_TESSERACT_NOT_FOUND_ERROR):
                logger.critical("Tesseract executable not found in PATH. OCR will fail for subsequent images too.")
                raise 
            logger.error(f"Error during OCR for image {i+1}/{len(image_objects)}: {e}", exc_info=True)
    
    full_ocr_text = "\n\n--- OCR Text from Image ---\n\n".join(ocr_text_parts)
    logger.info(f"OCR: Extracted {len(full_ocr_text)} chars from {images_ocrd} image(s) (out of {len(image_objects)} provided).")
    return full_ocr_text

# --- Stage 3: Text Cleaning and Normalization --- (Function as previously corrected)
def clean_and_normalize_text_content(text: str) -> str:
    if not text or not text.strip(): return ""
    logger.info(f"Starting text cleaning and normalization. Initial length: {len(text)}")
    text = re.sub(r'<script[^>]*>.*?</script>|<style[^>]*>.*?</style>|<[^>]+>', ' ', text, flags=re.I | re.S)
    text = re.sub(r'http\S+|www\S+|https\S+|\S*@\S*\s?', '', text, flags=re.MULTILINE)
    text = re.sub(r'[\n\t\r]+', ' ', text) 
    text = re.sub(r'\s+', ' ', text).strip() 
    text = re.sub(r'[^\w\s.,!?-]', '', text) 
    text_lower = text.lower()

    if not SPACY_MODEL_LOADED or not nlp_spacy_core:
        logger.warning("SpaCy model not loaded. Skipping tokenization/lemmatization. Returning regex-cleaned text.")
        return text_lower
    try:
        doc = nlp_spacy_core(text_lower, disable=['parser', 'ner']) 
        lemmatized_tokens = [token.lemma_ for token in doc if not token.is_stop and not token.is_punct and not token.is_space and len(token.lemma_) > 1]
        final_cleaned_text = " ".join(lemmatized_tokens)
        logger.info(f"SpaCy-based text cleaning and normalization complete. Final length: {len(final_cleaned_text)}")
        if final_cleaned_text: logger.info(f"Final cleaned text (first 200 chars): {final_cleaned_text[:200]}...")
        return final_cleaned_text
    except Exception as e:
        logger.error(f"SpaCy processing failed: {e}. Returning pre-SpaCy cleaned text.", exc_info=True)
        return text_lower

# --- Stage 4: Layout Reconstruction & Table Integration --- (Function as previously corrected)
def reconstruct_document_layout(text_content: str, tables_data: List[Any], file_type: str) -> str:
    if not text_content and not tables_data: return ""
    logger.info(f"Starting layout reconstruction for file type '{file_type}'. Initial text length: {len(text_content)}, Tables: {len(tables_data)}.")
    processed_text = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', text_content) 
    processed_text = re.sub(r'(\w+)-(\w+)', r'\1\2', processed_text)       

    if tables_data:
        table_md_parts = []
        for i, table_obj in enumerate(tables_data):
            header_md = f"\n\n[START OF TABLE {i+1}]\n"
            footer_md = f"\n[END OF TABLE {i+1}]\n"
            md_content = ""
            try:
                if PANDAS_AVAILABLE and pd and isinstance(table_obj, pd.DataFrame): 
                    md_content = table_obj.to_markdown(index=False)
                elif isinstance(table_obj, list) and table_obj and all(isinstance(r, list) for r in table_obj):
                    if table_obj and table_obj[0]: 
                        md_content = "| " + " | ".join(map(str, table_obj[0])) + " |\n"
                        md_content += "| " + " | ".join(["---"] * len(table_obj[0])) + " |\n"
                        for row_idx, row_data in enumerate(table_obj[1:]):
                            if len(row_data) == len(table_obj[0]): 
                                md_content += "| " + " | ".join(map(str, row_data)) + " |\n"
                            else:
                                logger.warning(f"Table {i+1}, row {row_idx+1} length mismatch with header. Skipping row.")
                    else: md_content = "[Empty Table Data]"
                else: md_content = str(table_obj)
            except Exception as e: 
                logger.warning(f"Table {i+1} to markdown conversion error: {e}. Using string representation."); 
                md_content = str(table_obj)
            if md_content: table_md_parts.append(header_md + md_content + footer_md)
        if table_md_parts: processed_text += "\n\n" + "\n\n".join(table_md_parts)
    
    processed_text = re.sub(r'\s+', ' ', processed_text).strip()
    logger.info(f"Layout reconstruction complete. Final text length: {len(processed_text)}")
    return processed_text

# --- Stage 5: Metadata Extraction --- (Function as previously corrected, uses original_file_name for logging)
def extract_document_metadata_info(file_path: str, processed_text: str, file_type_from_extraction: str, original_file_name: str, user_id: str) -> Dict[str, Any]:
    logger.info(f"Starting metadata extraction for: {original_file_name} (User: {user_id})")
    doc_meta = {'file_name': original_file_name, 'file_path_on_server': file_path, 'original_file_type': file_type_from_extraction,
                'processing_user': user_id, 'title': original_file_name, 'author': "Unknown", 'creation_date': None,
                'modification_date': None, 'page_count': 0, 'char_count_processed_text': len(processed_text),
                'named_entities': {}, 'structural_elements': "Paragraphs, Tables (inferred)", 'is_scanned_pdf': False}
    try:
        doc_meta['file_size_bytes'] = os.path.getsize(file_path)
        if PANDAS_AVAILABLE and pd:
            doc_meta['creation_date_os'] = pd.Timestamp(os.path.getctime(file_path), unit='s').isoformat()
            doc_meta['modification_date_os'] = pd.Timestamp(os.path.getmtime(file_path), unit='s').isoformat()
    except Exception as e: logger.warning(f"Metadata: OS metadata error for {original_file_name}: {e}")

    if file_type_from_extraction == '.pdf' and PYPDF2_AVAILABLE and PyPDF2:
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                info = reader.metadata
                if info:
                    if hasattr(info, 'title') and info.title: doc_meta['title'] = str(info.title).strip()
                    if hasattr(info, 'author') and info.author: doc_meta['author'] = str(info.author).strip()
                    if hasattr(info, 'creation_date') and info.creation_date and PANDAS_AVAILABLE and pd: 
                        doc_meta['creation_date'] = pd.Timestamp(info.creation_date).isoformat()
                doc_meta['page_count'] = len(reader.pages)
        except Exception as e: logger.warning(f"Metadata: PyPDF2 error for {original_file_name}: {e}", exc_info=True)
    elif file_type_from_extraction == '.docx' and DOCX_AVAILABLE and DocxDocument:
        try:
            doc = DocxDocument(file_path)
            props = doc.core_properties
            if props.title: doc_meta['title'] = props.title
            if props.author: doc_meta['author'] = props.author
            if props.created and PANDAS_AVAILABLE and pd: doc_meta['creation_date'] = pd.Timestamp(props.created).isoformat()
            doc_meta['page_count'] = sum(1 for p in doc.paragraphs if p.text.strip()) or 1
        except Exception as e: logger.warning(f"Metadata: DOCX error for {original_file_name}: {e}", exc_info=True)
    elif file_type_from_extraction == '.pptx' and PPTX_AVAILABLE and Presentation:
        try:
            prs = Presentation(file_path)
            props = prs.core_properties
            if props.title: doc_meta['title'] = props.title
            if props.author: doc_meta['author'] = props.author
            if props.created and PANDAS_AVAILABLE and pd: doc_meta['creation_date'] = pd.Timestamp(props.created).isoformat()
            doc_meta['page_count'] = len(prs.slides)
        except Exception as e: logger.warning(f"Metadata: PPTX error for {original_file_name}: {e}", exc_info=True)

    if doc_meta['page_count'] == 0 and processed_text: doc_meta['page_count'] = processed_text.count('\n\n') + 1

    if processed_text and SPACY_MODEL_LOADED and nlp_spacy_core:
        logger.info(f"Extracting named entities using SpaCy for {original_file_name}...")
        try:
            max_len = getattr(config, 'MAX_TEXT_LENGTH_FOR_NER', 500000) 
            text_for_ner = processed_text[:max_len]
            spacy_doc = nlp_spacy_core(text_for_ner) 
            ner_labels = nlp_spacy_core.pipe_labels.get("ner", [])
            entities = {label: [] for label in ner_labels}
            for ent in spacy_doc.ents:
                if ent.label_ in entities: entities[ent.label_].append(ent.text)
            doc_meta['named_entities'] = {k: list(set(v)) for k, v in entities.items() if v} 
            num_entities = sum(len(v_list) for v_list in doc_meta['named_entities'].values())
            logger.info(f"Extracted {num_entities} named entities for {original_file_name}.")
        except Exception as e: logger.error(f"Metadata: NER error for {original_file_name}: {e}", exc_info=True); doc_meta['named_entities'] = {}
    else: logger.info(f"Skipping NER for {original_file_name} (no text or SpaCy model not available/loaded)."); doc_meta['named_entities'] = {}
    
    logger.info(f"Metadata extraction complete for {original_file_name}.")
    return doc_meta

# --- Stage 6: Text Chunking ---
def chunk_document_into_segments(
    text_to_chunk: str,
    document_level_metadata: Dict[str, Any]
) -> List[Dict[str, Any]]:
    if not text_to_chunk or not text_to_chunk.strip():
        logger.warning(f"Chunking: No text for {document_level_metadata.get('file_name', 'unknown')}.")
        return []

    if not LANGCHAIN_SPLITTER_AVAILABLE or not RecursiveCharacterTextSplitter:
        logger.error("RecursiveCharacterTextSplitter not available. Cannot chunk text.")
        return []
        
    # CORRECTED: Use AI_CORE_CHUNK_SIZE and AI_CORE_CHUNK_OVERLAP from server/config.py
    chunk_s = getattr(config, 'AI_CORE_CHUNK_SIZE', 512) 
    chunk_o = getattr(config, 'AI_CORE_CHUNK_OVERLAP', 100) 

    original_doc_name_for_log = document_level_metadata.get('file_name', 'unknown')
    logger.info(f"Starting text chunking for {original_doc_name_for_log}. "
                f"Using config: CHUNK_SIZE={chunk_s}, CHUNK_OVERLAP={chunk_o}")
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_s,
        chunk_overlap=chunk_o,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""], 
        keep_separator=True 
    )

    try: raw_text_segments: List[str] = text_splitter.split_text(text_to_chunk)
    except Exception as e: 
        logger.error(f"Chunking: Error splitting text for {original_doc_name_for_log}: {e}", exc_info=True)
        return []
        
    output_chunks: List[Dict[str, Any]] = []
    base_file_name_for_ref = os.path.splitext(original_doc_name_for_log)[0] 

    for i, segment_content in enumerate(raw_text_segments):
        if not segment_content.strip(): 
            logger.debug(f"Skipping empty chunk at index {i} for {original_doc_name_for_log}.")
            continue

        chunk_specific_metadata = copy.deepcopy(document_level_metadata)
        
        qdrant_point_id = str(uuid.uuid4()) # Generate UUID for Qdrant

        chunk_specific_metadata['chunk_id'] = qdrant_point_id 
        chunk_specific_metadata['chunk_reference_name'] = f"{base_file_name_for_ref}_chunk_{i:04d}"
        chunk_specific_metadata['chunk_index'] = i
        chunk_specific_metadata['chunk_char_count'] = len(segment_content)
        
        output_chunks.append({
            'id': qdrant_point_id, 
            'text_content': segment_content,
            'metadata': chunk_specific_metadata 
        })
    
    logger.info(f"Chunking: Split '{original_doc_name_for_log}' into {len(output_chunks)} non-empty chunks with UUID IDs.")
    return output_chunks

# --- Stage 7: Embedding Generation --- (Function as previously corrected)
def generate_segment_embeddings(document_chunks: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    if not document_chunks: return []
    if not EMBEDDING_MODEL_LOADED or not document_embedding_model: # Check if model is loaded
        logger.error("Embedding model not loaded. Cannot generate embeddings.")
        for chunk_dict in document_chunks: chunk_dict['embedding'] = None # Ensure key exists
        return document_chunks

    # CORRECTED: Get model name from DOCUMENT_EMBEDDING_MODEL_NAME
    model_name_for_logging = getattr(config, 'DOCUMENT_EMBEDDING_MODEL_NAME', "Unknown Model")
    logger.info(f"Embedding: Starting embedding generation for {len(document_chunks)} chunks "
                f"using model: {model_name_for_logging}.")
    
    texts_to_embed: List[str] = []
    valid_chunk_indices: List[int] = []

    for i, chunk_dict in enumerate(document_chunks):
        text_content = chunk_dict.get('text_content')
        if text_content and text_content.strip():
            texts_to_embed.append(text_content)
            valid_chunk_indices.append(i)
        else:
            chunk_dict['embedding'] = None
            logger.debug(f"Embedding: Chunk {chunk_dict.get('id', i)} has no text content, skipping embedding.")

    if not texts_to_embed:
        logger.warning("Embedding: No actual text content found in any provided chunks to generate embeddings.")
        return document_chunks

    try:
        logger.info(f"Embedding: Generating embeddings for {len(texts_to_embed)} non-empty text segments...")
        embeddings_np_array = document_embedding_model.encode(texts_to_embed, show_progress_bar=False)
        
        for i, original_chunk_index in enumerate(valid_chunk_indices):
            if i < len(embeddings_np_array):
                document_chunks[original_chunk_index]['embedding'] = embeddings_np_array[i].tolist()
            else:
                logger.error(f"Embedding: Mismatch in embedding count for chunk at original index {original_chunk_index}. Assigning None.")
                document_chunks[original_chunk_index]['embedding'] = None
        
        logger.info(f"Embedding: Embeddings generated and assigned to {len(valid_chunk_indices)} chunks.")
        
    except Exception as e:
        logger.error(f"Embedding: Error during embedding generation with model {model_name_for_logging}: {e}", exc_info=True)
        for original_chunk_index in valid_chunk_indices:
            document_chunks[original_chunk_index]['embedding'] = None
        
    return document_chunks


# --- Main Orchestration Function (to be called by app.py) ---
def process_document_for_embeddings(file_path: str, original_name: str, user_id: str) -> List[Dict[str, Any]]:
    logger.info(f"ai_core: Orchestrating document processing for {original_name}, user {user_id}")
    if not os.path.exists(file_path): 
        logger.error(f"File not found at ai_core entry point: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    try:
        raw_content = extract_raw_content_from_file(file_path)
        
        ocr_text_output = ""
        if raw_content.get('is_scanned') and raw_content.get('images'):
            if PYTESSERACT_AVAILABLE and pytesseract:
                 ocr_text_output = perform_ocr_on_images(raw_content['images'])
            else:
                logger.warning(f"OCR requested for {original_name} but Pytesseract is not available/configured. Skipping OCR.")

        combined_text = raw_content.get('text_content', "")
        if ocr_text_output:
            if raw_content.get('is_scanned') and len(ocr_text_output) > len(combined_text) / 2:
                 combined_text = ocr_text_output + "\n\n" + combined_text
            else:
                 combined_text += "\n\n" + ocr_text_output
        
        if not combined_text.strip() and not raw_content.get('tables'):
            logger.warning(f"No text content for {original_name} after raw extraction/OCR, and no tables. Returning empty.")
            return []

        cleaned_text = clean_and_normalize_text_content(combined_text)
        if not cleaned_text.strip() and not raw_content.get('tables'):
            logger.warning(f"No meaningful text for {original_name} after cleaning, and no tables. Returning empty.")
            return []

        text_for_metadata_and_chunking = reconstruct_document_layout(
            cleaned_text,
            raw_content.get('tables', []),
            raw_content.get('file_type', '')
        )

        doc_metadata = extract_document_metadata_info(
            file_path,
            text_for_metadata_and_chunking, 
            raw_content.get('file_type', ''),
            original_name,
            user_id
        )
        doc_metadata['is_scanned_pdf'] = raw_content.get('is_scanned', False)

        # This now uses corrected config variable names and UUIDs for IDs
        chunks_with_metadata = chunk_document_into_segments( 
            text_for_metadata_and_chunking,
            doc_metadata
        )
        if not chunks_with_metadata:
            logger.warning(f"No chunks produced for {original_name}. Returning empty list.")
            return []

        # This now uses corrected config variable name for logging
        final_chunks_with_embeddings = generate_segment_embeddings(chunks_with_metadata)
        
        logger.info(f"ai_core: Successfully processed {original_name}. Generated {len(final_chunks_with_embeddings)} chunks.")
        return final_chunks_with_embeddings

    except Exception as e: 
        if _AI_CORE_TESSERACT_NOT_FOUND_ERROR and isinstance(e, _AI_CORE_TESSERACT_NOT_FOUND_ERROR):
            logger.critical(f"ai_core: Tesseract (OCR engine) was not found during processing of {original_name}. OCR could not be performed.", exc_info=False)
            raise 
        
        logger.error(f"ai_core: Critical error during processing of {original_name}: {e}", exc_info=True)
        raise